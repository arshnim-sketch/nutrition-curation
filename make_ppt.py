from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Bauhaus 색상
BLACK   = RGBColor(0x11, 0x11, 0x11)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xE6, 0x33, 0x29)
BLUE    = RGBColor(0x1B, 0x4F, 0xD8)
YELLOW  = RGBColor(0xF5, 0xC8, 0x00)
BG      = RGBColor(0xF5, 0xF0, 0xE8)
GRAY    = RGBColor(0x88, 0x88, 0x88)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def rgb(c): return c

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.line.width = line_w
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, x, y, w, h,
             size=Pt(16), bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Arial'
    return txb

def add_para(tf, text, size=Pt(14), bold=False, color=BLACK,
             space_before=Pt(4), bullet=False):
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.add_paragraph()
    p.space_before = space_before
    p.alignment = PP_ALIGN.LEFT
    if bullet:
        pPr = p._p.get_or_add_pPr()
        buNone = etree.SubElement(pPr, qn('a:buNone'))  # 기본 불릿 제거
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Arial'
    return p

# ─────────────────────────────────────────────
# SLIDE 0 — TITLE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 배경
add_rect(slide, 0, 0, W, H, fill=BLACK)

# 왼쪽 컬러 블록 3개 (바우하우스 기하학)
add_rect(slide, 0, 0, Inches(0.6), H, fill=RED)
add_rect(slide, Inches(0.6), 0, Inches(0.25), H, fill=YELLOW)
add_rect(slide, Inches(0.85), 0, Inches(0.12), H, fill=BLUE)

# 오른쪽 장식
add_rect(slide, W - Inches(1.2), Inches(0.5), Inches(1.2), Inches(1.2), fill=YELLOW)
add_rect(slide, W - Inches(0.8), Inches(2.0), Inches(0.8), Inches(0.8), fill=RED)
add_rect(slide, W - Inches(1.6), Inches(3.2), Inches(0.4), Inches(4.3), fill=BLUE)

# 타이틀
add_text(slide, 'AI 맞춤형\n영양제 큐레이션',
         Inches(1.4), Inches(1.2), Inches(9), Inches(2.8),
         size=Pt(56), bold=True, color=WHITE)

# 서브타이틀
add_text(slide, 'LLM + React 기반 개인 맞춤 임상 영양 큐레이션 서비스',
         Inches(1.4), Inches(4.2), Inches(9.5), Inches(0.8),
         size=Pt(20), bold=False, color=YELLOW)

# URL
add_text(slide, '07nutrition.vercel.app',
         Inches(1.4), Inches(5.3), Inches(6), Inches(0.5),
         size=Pt(14), bold=False, color=GRAY)

# ─────────────────────────────────────────────
# SLIDE 1 — 문제 정의
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=BG)
add_rect(slide, 0, 0, Inches(0.25), H, fill=RED)

# 슬라이드 번호 + 제목
add_text(slide, '01', Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.6),
         size=Pt(11), bold=True, color=RED)
add_text(slide, '프로젝트 개요 및 문제 정의',
         Inches(0.5), Inches(0.85), Inches(10), Inches(0.8),
         size=Pt(32), bold=True, color=BLACK)
add_rect(slide, Inches(0.5), Inches(1.6), Inches(4.5), Inches(0.06), fill=BLACK)

# 왼쪽 블록: 문제
add_rect(slide, Inches(0.5), Inches(1.9), Inches(5.6), Inches(4.8), fill=BLACK)
add_text(slide, '문제', Inches(0.7), Inches(2.0), Inches(2), Inches(0.5),
         size=Pt(11), bold=True, color=YELLOW)

txb = slide.shapes.add_textbox(Inches(0.7), Inches(2.55), Inches(5.2), Inches(4.0))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
add_para(tf, '수많은 영양제, 정말 나에게 맞는가?', size=Pt(18), bold=True, color=WHITE, space_before=Pt(2))
add_para(tf, '', size=Pt(6), color=WHITE)
add_para(tf, '— 마케팅 위주의 단순 추천 (증상 1:1 매핑)', size=Pt(14), color=GRAY, space_before=Pt(2))
add_para(tf, '— 영양소 과다 복용 위험 분석 부재', size=Pt(14), color=GRAY, space_before=Pt(2))
add_para(tf, '— 성분 간 흡수 방해 종합 분석 부재', size=Pt(14), color=GRAY, space_before=Pt(2))
add_para(tf, '  (예: 칼슘 ↔ 철분, 아연 ↔ 구리)', size=Pt(13), color=GRAY, space_before=Pt(1))

# 오른쪽 블록: 해결책
add_rect(slide, Inches(6.4), Inches(1.9), Inches(6.4), Inches(4.8), fill=BLUE)
add_text(slide, '해결책', Inches(6.6), Inches(2.0), Inches(3), Inches(0.5),
         size=Pt(11), bold=True, color=YELLOW)

txb2 = slide.shapes.add_textbox(Inches(6.6), Inches(2.55), Inches(6.0), Inches(4.0))
txb2.word_wrap = True
tf2 = txb2.text_frame
tf2.word_wrap = True
add_para(tf2, 'AI 임상 영양 큐레이션', size=Pt(18), bold=True, color=WHITE, space_before=Pt(2))
add_para(tf2, '', size=Pt(6), color=WHITE)
add_para(tf2, '— 나이·성별·증상 기반 맞춤 추천', size=Pt(14), color=WHITE, space_before=Pt(2))
add_para(tf2, '— KDRIs 기준 영양소 균형 자동 검증', size=Pt(14), color=WHITE, space_before=Pt(2))
add_para(tf2, '— 성분 간 시너지/독성 상호작용 분석', size=Pt(14), color=WHITE, space_before=Pt(2))
add_para(tf2, '— 복용 타이밍 가이드 제공', size=Pt(14), color=WHITE, space_before=Pt(2))

# ─────────────────────────────────────────────
# SLIDE 2 — 왜 LLM인가
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=BLACK)
add_rect(slide, 0, 0, Inches(0.25), H, fill=YELLOW)

add_text(slide, '02', Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.6),
         size=Pt(11), bold=True, color=YELLOW)
add_text(slide, '왜 LLM이 필요한가?',
         Inches(0.5), Inches(0.85), Inches(10), Inches(0.8),
         size=Pt(32), bold=True, color=WHITE)
add_rect(slide, Inches(0.5), Inches(1.6), Inches(3.5), Inches(0.06), fill=YELLOW)

reasons = [
    ('복잡한 다중 추론', BLUE,
     '수십 가지 성분 합산 → KDRIs 대조 →\n시너지/독성/타이밍 종합 판단\n\nif-else 룰베이스로는 구현 불가능'),
    ('자연어 설명 생성\n(XAI)', RED,
     '"왜 이 조합이 이 사람에게 맞는가"\n약사처럼 구체적이고 전문적인\n자연어 설명 자동 생성'),
    ('JSON 구조화\n강제 출력', YELLOW,
     'GPT-4o JSON Mode 활용\n프론트엔드가 즉시 렌더링 가능한\n예측 가능한 데이터 포맷 반환'),
]

for i, (title, color, body) in enumerate(reasons):
    x = Inches(0.5 + i * 4.2)
    add_rect(slide, x, Inches(2.0), Inches(3.9), Inches(4.8), fill=color)
    add_text(slide, title, x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.9),
             size=Pt(17), bold=True,
             color=BLACK if color == YELLOW else WHITE)
    add_rect(slide, x, Inches(3.0), Inches(3.9), Inches(0.04),
             fill=BLACK if color == YELLOW else WHITE)
    add_text(slide, body, x + Inches(0.2), Inches(3.1), Inches(3.55), Inches(3.5),
             size=Pt(14), bold=False,
             color=BLACK if color == YELLOW else WHITE)

# ─────────────────────────────────────────────
# SLIDE 3 — 아키텍처
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=BG)
add_rect(slide, 0, 0, Inches(0.25), H, fill=BLUE)

add_text(slide, '03', Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.6),
         size=Pt(11), bold=True, color=BLUE)
add_text(slide, '시스템 아키텍처',
         Inches(0.5), Inches(0.85), Inches(10), Inches(0.8),
         size=Pt(32), bold=True, color=BLACK)
add_rect(slide, Inches(0.5), Inches(1.6), Inches(4.0), Inches(0.06), fill=BLACK)

blocks = [
    (Inches(0.5),  Inches(1.9), Inches(3.0), Inches(4.8), BLACK, YELLOW, '🎨 UI Client',
     'React + Vite + TypeScript\nMobile-first 디자인\n2-Step 큐레이션 마법사\n실시간 체크박스 연동 UI'),
    (Inches(3.8),  Inches(1.9), Inches(3.0), Inches(4.8), BLUE,  WHITE,  '⚙️ AI System',
     'GPT-4o (JSON Mode)\n역할: 임상 영양사·약사\nRule-based Prompting\n멀티비타민 → 개별 성분 분해'),
    (Inches(7.1),  Inches(1.9), Inches(3.0), Inches(4.8), BLACK, RED,    '🔄 Data Pipeline',
     'Playwright 크롤링\n193개 실제 제품 카탈로그\nGoogle Sheets 실시간 로깅\nApps Script Webhook'),
    (Inches(10.4), Inches(1.9), Inches(2.4), Inches(4.8), YELLOW, BLACK, '☁️ Deploy',
     'GitHub\n+\nVercel\n(CI/CD 자동화)'),
]

for x, y, w, h, bg, tc, title, body in blocks:
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.15),
             w - Inches(0.3), Inches(0.55),
             size=Pt(14), bold=True, color=tc)
    add_rect(slide, x, y + Inches(0.75), w, Inches(0.04), fill=tc)
    add_text(slide, body, x + Inches(0.15), y + Inches(0.85),
             w - Inches(0.3), h - Inches(1.0),
             size=Pt(13), bold=False, color=tc)

# 화살표 텍스트
for ax in [Inches(3.55), Inches(6.85), Inches(10.15)]:
    add_text(slide, '→', ax, Inches(3.9), Inches(0.3), Inches(0.5),
             size=Pt(20), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 4 — 주요 기능
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=BLACK)
add_rect(slide, 0, 0, Inches(0.25), H, fill=RED)

add_text(slide, '04', Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.6),
         size=Pt(11), bold=True, color=RED)
add_text(slide, 'Key Highlights',
         Inches(0.5), Inches(0.85), Inches(10), Inches(0.8),
         size=Pt(32), bold=True, color=WHITE)
add_rect(slide, Inches(0.5), Inches(1.6), Inches(3.5), Inches(0.06), fill=RED)

features = [
    (RED,    '01', '2-Step 큐레이션 마법사',
     '관심 건강 영역 선택 →\n세부 증상 선택\n인지 부하 최소화 UX'),
    (BLUE,   '02', 'Interactive Nutrient Balance',
     'KDRIs 대비 영양소 막대 시각화\n체크박스 OFF 시 해당 영양소\n막대 실시간 감소 애니메이션'),
    (YELLOW, '03', 'Clinical Interactions',
     '3가지 유형 자동 분류:\n시너지(↑) / 주의(✕) /\n복용 타이밍(◷)'),
]

for i, (color, num, title, body) in enumerate(features):
    x = Inches(0.5 + i * 4.2)
    # 번호 블록
    add_rect(slide, x, Inches(2.0), Inches(0.65), Inches(0.65), fill=color)
    add_text(slide, num, x, Inches(2.0), Inches(0.65), Inches(0.65),
             size=Pt(16), bold=True,
             color=BLACK if color == YELLOW else WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.8), Inches(2.05), Inches(3.1), Inches(0.7),
             size=Pt(16), bold=True, color=color)
    add_rect(slide, x, Inches(2.75), Inches(3.9), Inches(0.04), fill=color)
    add_text(slide, body, x, Inches(2.9), Inches(3.9), Inches(3.0),
             size=Pt(14), color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 5 — 한계 및 개선
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=BG)
add_rect(slide, 0, 0, Inches(0.25), H, fill=YELLOW)

add_text(slide, '05', Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.6),
         size=Pt(11), bold=True, color=BLACK)
add_text(slide, '한계점 및 향후 개선 방향',
         Inches(0.5), Inches(0.85), Inches(11), Inches(0.8),
         size=Pt(32), bold=True, color=BLACK)
add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.0), Inches(0.06), fill=BLACK)

limits = [
    ('비용 & 레이턴시', RED,
     '현재: LLM 매 호출마다 수 초 대기\n→ 향후: 유사 증상 조합 결과를 Vector DB에 캐싱,\n   즉시 리턴 아키텍처로 개선'),
    ('보안 & 백엔드', BLUE,
     '현재: 프론트 직접 API 호출 (키 노출 위험)\n→ 향후: FastAPI/Node 백엔드 + LangChain/\n   LangGraph 엔드포인트로 이전'),
    ('환각 관리\n(Hallucination)', BLACK,
     '현재: DB 내 제품만 추천하도록 프롬프트 제어\n+ 클라이언트 단 텍스트 매칭 검증\n→ 향후: Fine-tuning 또는 Guard Rail 적용'),
]

for i, (title, color, body) in enumerate(limits):
    y = Inches(2.0 + i * 1.7)
    add_rect(slide, Inches(0.5), y, Inches(0.25), Inches(1.4), fill=color)
    add_text(slide, title, Inches(0.95), y + Inches(0.1), Inches(2.5), Inches(0.7),
             size=Pt(15), bold=True, color=color)
    add_text(slide, body, Inches(0.95), y + Inches(0.65), Inches(11.8), Inches(0.9),
             size=Pt(13), color=BLACK)

# ─────────────────────────────────────────────
# SLIDE 6 — CLOSING
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=BLACK)

# 바우하우스 기하학
add_rect(slide, 0, 0, Inches(4.5), H, fill=RED)
add_rect(slide, Inches(4.5), 0, Inches(0.25), H, fill=YELLOW)
add_rect(slide, Inches(4.75), 0, Inches(0.12), H, fill=BLUE)

add_text(slide, 'Thank you',
         Inches(5.2), Inches(2.4), Inches(7.5), Inches(1.4),
         size=Pt(52), bold=True, color=WHITE)
add_text(slide, '07nutrition.vercel.app',
         Inches(5.2), Inches(4.0), Inches(7.5), Inches(0.6),
         size=Pt(18), color=YELLOW)
add_text(slide, 'github.com/arshnim-sketch/nutrition-curation',
         Inches(5.2), Inches(4.6), Inches(7.5), Inches(0.5),
         size=Pt(13), color=GRAY)

# 왼쪽 텍스트 (세로 on red)
add_text(slide, 'AI NUTRITION\nCURATION',
         Inches(0.3), Inches(2.2), Inches(4.0), Inches(3.0),
         size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
prs.save('/Users/miyo/07_nutrition/AI_영양제큐레이션_발표자료.pptx')
print('✓ 저장 완료: AI_영양제큐레이션_발표자료.pptx')
